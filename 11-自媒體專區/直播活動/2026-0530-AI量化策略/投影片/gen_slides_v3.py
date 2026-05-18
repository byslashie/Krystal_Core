#!/usr/bin/env python3
"""
gen_slides_v3.py — slides_20260530_v3.pptx
按照 slide-spec.html 品牌規格重新生成
- 亮色系版：s-divider 全改為暖米白 GRAY_50
- 只有封面 (slide 01) + 結尾 (slide 19) 使用深色 NAVY_900
- 字體：Calibri Bold（Outfit 替代）+ Calibri（Inter）+ Consolas（Mono）
- 顏色：完整依照品牌 token
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 品牌 Color Tokens（完整對應 slide-spec.html）─────────────────
ORANGE_100  = RGBColor(255, 234, 214)   # --orange-100  (hint bg)
ORANGE_200  = RGBColor(245, 204, 166)   # --orange-200
ORANGE_300  = RGBColor(232, 171, 120)   # --orange-300  (dark-bg label)
ORANGE_400  = RGBColor(232, 141,  91)   # --orange-400
ORANGE_500  = RGBColor(232, 115,  74)   # --orange-500  ★主色 CTA/Tag
ORANGE_600  = RGBColor(204,  91,  51)   # --orange-600  (label on light)
ORANGE_700  = RGBColor(170,  68,  35)   # --orange-700
NAVY_700    = RGBColor( 23,  46,  89)   # --navy-700
NAVY_900    = RGBColor( 14,  24,  51)   # --navy-900    ★封面底色
CORAL       = RGBColor(255, 107, 157)   # --coral       (Aurora accent)
CYAN        = RGBColor( 77, 208, 225)   # --cyan
GOLD        = RGBColor(212, 165, 116)   # --gold
GRAY_50     = RGBColor(250, 247, 242)   # --gray-50     ★內容頁底色 (暖米白)
GRAY_100    = RGBColor(242, 237, 230)   # --gray-100
GRAY_200    = RGBColor(230, 222, 210)   # --gray-200
GRAY_400    = RGBColor(180, 165, 145)   # --gray-400
GRAY_500    = RGBColor(140, 127, 112)   # --gray-500    text-mute
GRAY_700    = RGBColor( 70,  60,  50)   # --gray-700    text-base
GRAY_900    = RGBColor( 19,  20,  30)   # --gray-900    text-strong
POS         = RGBColor(193,  59,  44)   # --pos  (上漲紅)
WARN        = RGBColor(245, 158,  11)   # --warn (警示黃)
INFO        = RGBColor( 59, 130, 246)   # --info
WHITE       = RGBColor(255, 255, 255)
CARD_BG     = RGBColor(255, 255, 255)   # 卡片白底
DIVIDER_LINE= RGBColor(220, 213, 204)   # 細分隔線

# ── 字體（PPT 安全替代）────────────────────────────────────────
F_DISPLAY = "Calibri"    # Outfit → Calibri Bold
F_BODY    = "Calibri"    # Inter  → Calibri
F_MONO    = "Consolas"   # JetBrains Mono → Consolas

# ── 畫布（20×11.25 inch = 1920×1080 等效）─────────────────────
W  = Inches(20)
H  = Inches(11.25)
MX = Inches(1.396)   # 左右安全邊距 134px
MY = Inches(0.667)   # 上下安全邊距 64px

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── 工具函數 ───────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0):
    from pptx.util import Pt as P2
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = P2(lw)
    else:
        s.line.fill.background()
    return s

def tb(slide, x, y, w, h, text, pt, color,
       bold=False, italic=False, font=F_BODY,
       align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box

def watermark(slide):
    tb(slide, 17.1, 10.78, 2.7, 0.32,
       "斜槓姐姐 · @byslashie", 10, GRAY_400, font=F_MONO)

def pagenum(slide, n, total=19):
    tb(slide, 1.396, 10.78, 1.8, 0.32,
       f"{n:02d} / {total:02d}", 10, GRAY_400, font=F_MONO)

def footer_social(slide):
    tb(slide, 1.396, 10.78, 4, 0.32,
       "byslashie.com", 11, GRAY_500, font=F_MONO)
    tb(slide, 15.5, 10.78, 4, 0.32,
       "Threads @byslashie", 11, GRAY_500, font=F_MONO, align=PP_ALIGN.RIGHT)

def accent_line(slide, x=1.396, y=1.3):
    rect(slide, x, y, 0.83, 0.042, fill=ORANGE_500)

def tag(slide, x, y, text, color=None, dark_bg=False):
    c = color or (ORANGE_300 if dark_bg else ORANGE_600)
    tb(slide, x, y, 14, 0.32, text, 13, c, bold=True,
       font=F_DISPLAY)

def title(slide, x, y, text, pt=42, color=None, bold=True):
    c = color or GRAY_900
    tb(slide, x, y, 17.2, 0.85, text, pt, c, bold=bold,
       font=F_DISPLAY)

def bullet(slide, x, y, text, col=None, size=19, dot=None):
    dc = dot or ORANGE_500; tc = col or GRAY_700
    rect(slide, x, y+0.1, 0.1, 0.1, fill=dc)
    tb(slide, x+0.28, y, 16.5, 0.48, text, size, tc, wrap=True)

def card(slide, x, y, w, h, top_color=ORANGE_500):
    rect(slide, x, y, w, h, fill=CARD_BG, line=DIVIDER_LINE, lw=0.75)
    rect(slide, x, y, w, 0.055, fill=top_color)

# ── s-divider 版型（按 spec：亮色 GRAY_50 + 橙大字編號）─────────
def section_slide(n_str, title_line1, title_line2, sub):
    sl = prs.slides.add_slide(BLANK)
    bg(sl, GRAY_50)
    # 右下暈光裝飾（用橙色矩形低透明度模擬）
    rect(sl, 12, 6, 8, 5.5, fill=ORANGE_100)   # 右下橙暈
    rect(sl, 0,  0, 6, 4.5, fill=RGBColor(232,248,252))  # 左上青暈
    bg(sl, GRAY_50)  # 重設底色蓋住暈光邊緣（暈光改用z-order疊加）
    # 重畫：暈光改為透明矩形 → pptx 無法做透明，改用接近色塊
    # 實際效果：左上淡青、右下淡橙作為氛圍色
    # 章節編號（大橙字）
    tb(sl, 1.396, 2.8, 4, 1.5, n_str, 88, ORANGE_500,
       bold=False, font=F_MONO)
    # 主標題
    tb(sl, 1.396, 4.5, 15, 1.1, title_line1, 64, GRAY_900,
       bold=True, font=F_DISPLAY)
    if title_line2:
        tb(sl, 1.396, 5.7, 15, 1.1, title_line2, 64, GRAY_900,
           bold=True, font=F_DISPLAY)
    # sub
    tb(sl, 1.396, 7.2, 15, 0.5, sub, 19, GRAY_500,
       italic=True, font=F_BODY)
    # accent bar
    rect(sl, 1.396, 7.95, 0.83, 0.055, fill=ORANGE_500)
    watermark(sl)
    return sl

# ═══════════════════════════════════════════════════════════════
# SLIDE 01 — 封面（深色 s-cover）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, NAVY_900)

# Aurora 裝飾色塊（低飽和模擬漸層）
rect(sl, 0, 0, 8, 11.25, fill=RGBColor(20, 36, 72))    # 左側略亮 navy
rect(sl, 12, 0, 8, 6,    fill=RGBColor(18, 30, 58))    # 右上

# Top bar line
rect(sl, 1.396, 1.8, 0.83, 0.055, fill=ORANGE_500)
tb(sl, 1.396, 2.0, 15, 0.4,
   "AI 策略入門班  ·  2026.05.30", 14, ORANGE_300,
   bold=True, font=F_DISPLAY)

# 主標題
tb(sl, 1.396, 3.0, 16, 1.2, "不會寫程式", 80, WHITE,
   bold=True, font=F_DISPLAY)
tb(sl, 1.396, 4.3, 16, 1.2, "也能做策略", 80, ORANGE_300,
   bold=True, font=F_DISPLAY)

# 副標
tb(sl, 1.396, 5.8, 15, 0.6,
   "用 AI 打造你的 XQ 投資助手", 22,
   RGBColor(200, 195, 210), font=F_BODY)

# Aurora 腰帶（模擬珊瑚/青色光帶）
rect(sl, 0, 8.2, 20, 0.06, fill=ORANGE_500)

footer_social(sl)
watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 02 — 今天議程（s-content 亮色）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "TODAY'S AGENDA")
title(sl, 1.396, 1.9, "今天你會學到什麼")

items = [
    "以前建構策略的流程與痛點",
    "用 Perplexity + 論文找到有根據的因子",
    "Claude Skill 一步步釐清目標、建立 XS 策略",
    "XQ 回測 → 結果存進 HTML 視覺化分析",
]
for i, it in enumerate(items):
    y = 3.1 + i * 0.9
    rect(sl, 1.396, y+0.12, 0.38, 0.38, fill=ORANGE_500)
    tb(sl, 1.396, y+0.13, 0.38, 0.38,
       str(i+1), 15, WHITE, bold=True,
       font=F_DISPLAY, align=PP_ALIGN.CENTER)
    tb(sl, 1.95, y, 16.0, 0.48, it, 20, GRAY_700)

pagenum(sl, 2); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 03 — 以前怎麼做（s-compare 對比）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "BEFORE · 舊方法")
title(sl, 1.396, 1.9, "以前建構策略怎麼做？")

# 流程箭頭框
flow_row1 = ["有想法", "找指標", "憑感覺回測"]
flow_row2 = ["修改策略", "實單交易", "結果曝虧"]
bw, bh, gap_x = 3.6, 0.68, 1.1
base_x = 1.396
y1 = 3.5; y2 = 5.3

for i, label in enumerate(flow_row1):
    x = base_x + i * (bw + gap_x)
    card(sl, x, y1, bw, bh, GRAY_400)
    tb(sl, x, y1+0.14, bw, bh, label, 20, GRAY_900,
       bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    if i < 2:
        rect(sl, x+bw+0.1, y1+0.3, gap_x-0.15, 0.08, fill=ORANGE_500)

for i, label in enumerate(flow_row2):
    x = base_x + (2-i) * (bw + gap_x)
    card(sl, x, y2, bw, bh, POS if label == "結果曝虧" else GRAY_400)
    tb(sl, x, y2+0.14, bw, bh, label, 20, GRAY_900,
       bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    if i < 2:
        rect(sl, x-(gap_x-0.15+0.1), y2+0.3, gap_x-0.15, 0.08, fill=ORANGE_500)

# 下方評語
rect(sl, 1.396, 7.5, 17.2, 0.75, fill=RGBColor(255, 235, 230), line=POS, lw=1.5)
tb(sl, 1.7, 7.62, 16.5, 0.5,
   "耗時費力 · 靠感覺 · 沒有學術根據 · 勝率無法改善",
   17, POS, bold=True)

pagenum(sl, 3); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 04 — 現在的流程（s-content）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "NOW WITH AI · 新流程")
title(sl, 1.396, 1.9, "加入 AI 之後")

steps_ai = [
    ("有想法", ORANGE_500),
    ("Claude Skill", ORANGE_500),
    ("Perplexity 找論文", ORANGE_500),
    ("Claude 寫 XS", ORANGE_500),
    ("XQ 回測", ORANGE_500),
    ("AI 解讀報告", ORANGE_500),
]
bw2 = (17.2 - 0.2 * (len(steps_ai)-1)) / len(steps_ai)
y_row = 3.2
for i, (label, c) in enumerate(steps_ai):
    x = 1.396 + i * (bw2 + 0.2)
    card(sl, x, y_row, bw2, 0.72, c)
    tb(sl, x, y_row+0.16, bw2, 0.72, label, 15, NAVY_900,
       bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    if i < len(steps_ai)-1:
        tb(sl, x+bw2+0.02, y_row+0.22, 0.2, 0.4, "▶", 14, ORANGE_500, align=PP_ALIGN.CENTER)

# 第二行 — 循環改進
improvements = [
    "每個環節 AI 輔助",
    "有論文根據，策略有邏輯",
    "省去 80% 試錯時間",
]
for i, txt in enumerate(improvements):
    y = 4.6 + i * 0.85
    rect(sl, 1.396, y+0.16, 0.1, 0.1, fill=ORANGE_500)
    tb(sl, 1.75, y, 16.0, 0.48, txt, 20, GRAY_700)

# 強調框
rect(sl, 1.396, 8.2, 17.2, 0.8, fill=ORANGE_100, line=ORANGE_500, lw=1.5)
tb(sl, 1.7, 8.35, 16.5, 0.5,
   "今天直播：我帶你把這個流程跑一遍", 20, ORANGE_600,
   bold=True, font=F_DISPLAY)

pagenum(sl, 4); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 05 — 章節 01（s-divider 亮色）
# ═══════════════════════════════════════════════════════════════
section_slide("01", "為什麼從論文", "找因子？", "學術依據 vs 網路資訊")

# ═══════════════════════════════════════════════════════════════
# SLIDE 06 — 論文 vs 網路（s-compare）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "THE PROBLEM")
title(sl, 1.396, 1.9, "為什麼要找論文？")

# Left — 網路資訊
lx, ly, lw, lh = 1.396, 3.2, 7.8, 5.8
rect(sl, lx, ly, lw, lh, fill=RGBColor(244,242,248), line=RGBColor(180,170,200), lw=1.5)
rect(sl, lx, ly, lw, 0.055, fill=NAVY_700)
tb(sl, lx+0.28, ly+0.2, lw-0.4, 0.6, "❌  網路資訊", 24, NAVY_700, bold=True, font=F_DISPLAY)
for i, t in enumerate(["來源不明、難以驗證", "容易過時、流傳錯誤", "不知道有沒有真的效果"]):
    bullet(sl, lx+0.28, ly+1.15+i*1.1, t, col=GRAY_700, dot=RGBColor(180,60,60))

# Right — 學術論文
rx, ry = 10.6, 3.2
rect(sl, rx, ry, lw, lh, fill=RGBColor(255,245,235), line=ORANGE_500, lw=1.5)
rect(sl, rx, ry, lw, 0.055, fill=ORANGE_500)
tb(sl, rx+0.28, ry+0.2, lw-0.4, 0.6, "✅  學術論文", 24, NAVY_700, bold=True, font=F_DISPLAY)
for i, t in enumerate(["有方法論、可重現驗證", "清楚的因子條件", "建立策略有根有據"]):
    bullet(sl, rx+0.28, ry+1.15+i*1.1, t, col=GRAY_700, dot=ORANGE_500)

pagenum(sl, 6); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 07 — 章節 02（s-divider 亮色）
# ═══════════════════════════════════════════════════════════════
section_slide("02", "找論文", "兩種方法", "免費版 & 進階版")

# ═══════════════════════════════════════════════════════════════
# SLIDE 08 — Perplexity（s-content）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "SIMPLE · 免費 · 零設定")
title(sl, 1.396, 1.9, "Perplexity Academic 模式")

for i, (n, step) in enumerate([
    ("1", "開啟 perplexity.ai → 右上角切換 Academic 模式"),
    ("2", "輸入英文關鍵字，例：momentum factor Taiwan stock"),
    ("3", "把搜尋摘要貼給 Claude，請它分析因子條件"),
]):
    y = 3.2 + i * 1.15
    rect(sl, 1.396, y, 0.52, 0.52, fill=ORANGE_500)
    tb(sl, 1.396, y+0.06, 0.52, 0.52, n, 20, WHITE,
       bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    tb(sl, 2.1, y+0.04, 16.3, 0.5, step, 20, GRAY_700)

# 關鍵字示範框
rect(sl, 1.396, 7.0, 17.2, 1.8, fill=NAVY_900)
rect(sl, 1.396, 7.0, 0.083, 1.8, fill=ORANGE_500)
kw_items = [
    '"momentum factor Taiwan stock"',
    '"quality factor ROE profitability returns"',
    '"Fama French three factor model"',
]
for i, kw in enumerate(kw_items):
    tb(sl, 1.65, 7.1+i*0.5, 16.5, 0.45, kw, 16, ORANGE_300, font=F_MONO)

# 底部提示
rect(sl, 1.396, 9.15, 17.2, 0.72, fill=ORANGE_100, line=ORANGE_500, lw=1.5)
tb(sl, 1.7, 9.28, 16.5, 0.5,
   "完全免費，人人可用，今天 DEMO 用這個", 19, ORANGE_600, bold=True)

pagenum(sl, 8); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 09 — MCP 進階版（s-twocol）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "ADVANCED · 我自己用的方式")
title(sl, 1.396, 1.9, "MCP：Consensus + Scholar Gateway")

tools = [
    ("Consensus", "找實證結論", "輸入研究問題\n→ 直接輸出有/無支持", ORANGE_500),
    ("Scholar Gateway", "語意搜尋", "找最新文獻\n→ 多角度交叉驗證", NAVY_700),
]
for i, (name, sub1, sub2, top) in enumerate(tools):
    cx = 1.396 + i * 9.1
    card(sl, cx, 3.2, 8.6, 5.0, top)
    tb(sl, cx+0.35, 3.45, 7.8, 0.7, name, 32, NAVY_900, bold=True, font=F_DISPLAY)
    tb(sl, cx+0.35, 4.3,  7.8, 0.4, sub1, 19, ORANGE_600, bold=True, font=F_DISPLAY)
    tb(sl, cx+0.35, 4.8,  7.8, 0.9, sub2, 17, GRAY_700, wrap=True)

rect(sl, 1.396, 9.0, 17.2, 0.72, fill=GRAY_100, line=GRAY_200, lw=1)
tb(sl, 1.7, 9.12, 16.5, 0.5,
   "需設定 MCP Server，適合進階使用者；今天先用 Perplexity",
   16, GRAY_500, italic=True)

pagenum(sl, 9); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — 章節 03（s-divider 亮色）
# ═══════════════════════════════════════════════════════════════
section_slide("03", "Claude Skill", "策略建立向導", "一步步釐清你的策略目標")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Skill 介紹（s-content）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "THE SKILL")
title(sl, 1.396, 1.9, "什麼是 Claude Skill？")

tb(sl, 1.396, 2.9, 17.2, 0.9,
   "一段 Prompt，貼進 Claude 對話框就能用。\n不需要 Claude Pro、不需要設定，開新對話貼上就是你的策略顧問。",
   17, GRAY_700, wrap=True)

for i, (n, step) in enumerate([
    ("1", "打開 claude.ai → 開新對話"),
    ("2", "把 Skill Prompt 全選複製，貼進去"),
    ("3", "Claude 會問你策略方向 → 一步步建立邏輯"),
]):
    y = 4.3 + i * 0.9
    rect(sl, 1.396, y+0.1, 0.38, 0.38, fill=ORANGE_500)
    tb(sl, 1.396, y+0.12, 0.38, 0.38, n, 16, WHITE,
       bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    tb(sl, 1.95, y+0.04, 16.0, 0.45, step, 20, GRAY_700)

# CTA 橙色強調
rect(sl, 1.396, 8.2, 17.2, 0.9, fill=ORANGE_500)
tb(sl, 0, 8.35, 20, 0.6,
   "今天直播結束後，填表就送你這個 Skill！", 24,
   WHITE, bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)

pagenum(sl, 11); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Skill 五步驟（s-stat）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "HOW IT WORKS")
title(sl, 1.396, 1.9, "Skill 一步步問你什麼")

questions = [
    "商品類型（台股個股 / ETF / 美股）",
    "操作週期（日線 / 週線 / 月線）",
    "策略邏輯（順勢 / 逆勢 / 基本面 / 籌碼）",
    "持有時間（短線 / 波段 / 長持）",
    "風控需求（停損幅度、每月操作次數）",
]
for i, q in enumerate(questions):
    y = 3.0 + i * 1.0
    rect(sl, 1.396, y, 0.5, 0.5, fill=ORANGE_500)
    tb(sl, 1.396, y+0.06, 0.5, 0.5,
       str(i+1), 18, WHITE, bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    tb(sl, 2.1, y+0.06, 16.0, 0.46, q, 19, GRAY_700)

tb(sl, 1.396, 9.1, 17.2, 0.55,
   "→ 收到答案後，Skill 輸出 XS 策略框架（基本面 + 籌碼面 + 技術面三層）",
   17, ORANGE_600, bold=True)

pagenum(sl, 12); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — 章節 04（s-divider 亮色）
# ═══════════════════════════════════════════════════════════════
section_slide("04", "XQ 實作", "+ 回測", "把 AI 輸出的 XS 框架跑起來")

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — XQ 四步驟（s-stat 四卡片）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "LIVE DEMO")
title(sl, 1.396, 1.9, "XQ 實作四步驟")

grid = [
    ("01", "選股中心\n設定三因子條件",  1.396, 3.2),
    ("02", "建立 XS 交易腳本\n（AI 生成）",  10.3,  3.2),
    ("03", "設計出場邏輯\n（季度換倉）", 1.396, 6.6),
    ("04", "執行回測\n分析報告",         10.3,  6.6),
]
for num, label, cx, cy in grid:
    card(sl, cx, cy, 7.8, 3.1, ORANGE_500)
    tb(sl, cx+0.3, cy+0.25, 2.5, 1.1, num, 60, ORANGE_400,
       bold=False, font=F_MONO)
    tb(sl, cx+0.3, cy+1.5, 7.0, 0.9, label, 22, GRAY_900,
       bold=True, font=F_DISPLAY, wrap=True)

pagenum(sl, 14); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — 回測結果（s-stat）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "RESULT · 因子 ETF 策略")
title(sl, 1.396, 1.9, "5 年回測成績")

stat_data = [
    ("5年最大報酬", "+112.42%", "優化版 · 3 因子", POS, ORANGE_500),
    ("基礎版報酬",  "+86.42%",  "3 因子基礎設定",  WARN, WARN),
    ("0050 同期",   "+62%",     "大盤對比基準",    GRAY_500, GRAY_400),
]
for i, (lab, num, delta, numcol, topcol) in enumerate(stat_data):
    cx = 1.396 + i * 6.0
    card(sl, cx, 3.2, 5.6, 4.2, topcol)
    tb(sl, cx+0.3, 3.55, 5.0, 0.45, lab, 17, GRAY_700)
    tb(sl, cx+0.3, 4.1,  5.0, 1.2,  num, 54, numcol, bold=True, font=F_DISPLAY)
    tb(sl, cx+0.3, 6.8,  5.0, 0.4,  delta, 15, GRAY_400)

# 附注
tb(sl, 1.396, 8.5, 17.2, 0.5,
   "2022 年下半年起逆轉領先大盤，2019–2022 年表現弱於 0050 / 詳見 byslashie.com",
   15, GRAY_400, italic=True)

pagenum(sl, 15); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — 存進 HTML（s-content）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "ANALYSIS")
title(sl, 1.396, 1.9, "回測結果存進 HTML 分析")

for i, (n, step) in enumerate([
    ("1", "XQ 回測完成 → 匯出 CSV"),
    ("2", "把 CSV 拖入 Demo 分析系統"),
    ("3", "AI 自動計算 CAGR / Sharpe / MDD 等 18 項 KPI"),
]):
    y = 3.2 + i * 1.15
    rect(sl, 1.396, y+0.1, 0.38, 0.38, fill=ORANGE_500)
    tb(sl, 1.396, y+0.12, 0.38, 0.38, n, 16, WHITE,
       bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    tb(sl, 1.95, y+0.04, 16.0, 0.5, step, 20, GRAY_700)

# 右側 KPI 卡片
kpi_items = ["+CAGR", "Sharpe", "MaxDD", "WinRate"]
for i, k in enumerate(kpi_items):
    cx = 1.396 + i * 4.4
    rect(sl, cx, 7.0, 4.0, 1.5, fill=ORANGE_100, line=ORANGE_500, lw=1)
    rect(sl, cx, 7.0, 4.0, 0.055, fill=ORANGE_500)
    tb(sl, cx, 7.3, 4.0, 0.8, k, 28, ORANGE_600,
       bold=True, font=F_MONO, align=PP_ALIGN.CENTER)

pagenum(sl, 16); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — AI 解讀回測（s-quote + code block）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, GRAY_50)

accent_line(sl)
tag(sl, 1.396, 1.5, "AI REVIEW")
title(sl, 1.396, 1.9, "讓 AI 幫你解讀回測報告")

# 深色 code block
rect(sl, 1.396, 3.0, 17.2, 5.2, fill=NAVY_900)
rect(sl, 1.396, 3.0, 0.083, 5.2, fill=ORANGE_500)   # 左側橙線
code_lines = [
    "這是我的回測結果：[貼上數據]",
    "",
    "請幫我分析：",
    "1. 勝率和賠率是否合理？",
    "2. 這個策略的弱點在哪？",
    "3. 給我 3 個具體改進方向",
]
for i, line in enumerate(code_lines):
    color = ORANGE_300 if line.startswith("請") else (
            RGBColor(130, 220, 180) if line.startswith(("1", "2", "3")) else
            RGBColor(200, 195, 215))
    tb(sl, 1.65, 3.25+i*0.72, 16.5, 0.65, line, 17,
       color, font=F_MONO)

# 說明
rect(sl, 1.396, 8.55, 17.2, 0.72, fill=ORANGE_100, line=ORANGE_500, lw=1.5)
tb(sl, 1.7, 8.68, 16.5, 0.5,
   "Claude 會給你策略診斷 + 3 個具體優化方向", 19, ORANGE_600, bold=True)

pagenum(sl, 17); watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — 禮物章節（深色 s-closing style）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, NAVY_900)

# 裝飾色塊
rect(sl, 0,  0, 20, 11.25, fill=NAVY_900)
rect(sl, 6,  2,  8,  7,    fill=RGBColor(20, 38, 78))  # 中央亮圓

tb(sl, 0, 2.6, 20, 1.5, "🎁", 80, WHITE, align=PP_ALIGN.CENTER)
tb(sl, 0, 4.2, 20, 1.1, "今天的禮物", 64, WHITE,
   bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
tb(sl, 0, 5.5, 20, 0.6,
   "帶走這個工具，自己找下一個因子",
   22, ORANGE_300, font=F_BODY, align=PP_ALIGN.CENTER)

watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — 禮物包 + 結尾（深色 s-closing）
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, NAVY_900)

# 頂部橙線
rect(sl, 1.396, 1.7, 0.83, 0.055, fill=ORANGE_500)
tag(sl, 1.396, 1.9, "GIFT  ·  填表免費領取", dark_bg=True)

tb(sl, 1.396, 2.5, 17.2, 1.0,
   "禮物包內容", 50, WHITE, bold=True, font=F_DISPLAY)

gifts = [
    "📄  Claude Skill Prompt（貼入 Claude 即用，免費）",
    "🔍  Perplexity 學術搜尋關鍵字清單",
    "📊  XS 策略框架模板（三層選股邏輯可直接修改）",
]
for i, g in enumerate(gifts):
    y = 3.8 + i * 1.1
    rect(sl, 1.396, y+0.12, 0.42, 0.42, fill=ORANGE_500)
    tb(sl, 1.396, y+0.12, 0.42, 0.42, str(i+1), 17, WHITE,
       bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    tb(sl, 2.0, y, 16.8, 0.52, g, 21, RGBColor(235, 228, 220))

# CTA 橙帶
rect(sl, 1.396, 8.4, 17.2, 0.92, fill=ORANGE_500)
tb(sl, 0, 8.58, 20, 0.6,
   "直播留言「積木」→ 填表領取三份禮物",
   25, WHITE, bold=True, font=F_DISPLAY, align=PP_ALIGN.CENTER)

footer_social(sl)
watermark(sl)

# ═══════════════════════════════════════════════════════════════
# 儲存
# ═══════════════════════════════════════════════════════════════
out = "/Users/chenyu/Desktop/Projects/Krystal_Core/07-工具/slides_20260530_v3.pptx"
prs.save(out)
print("✅ slides_20260530_v3.pptx 已生成（19頁）")
print("   封面/結尾：深色 NAVY_900")
print("   章節頁(05/07/10/13/18)：暖米白 GRAY_50（亮色系）")
print("   內容頁：暖米白 GRAY_50 + 橙色主色")
