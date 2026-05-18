#!/usr/bin/env python3
"""
Generate slides_20260530_v2.pptx - AI 策略入門班 19-slide deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Color palette ──────────────────────────────────────────────
ORANGE_300 = RGBColor(232, 171, 120)
ORANGE_500 = RGBColor(232, 115,  74)
ORANGE_600 = RGBColor(204,  91,  51)
NAVY_700   = RGBColor( 23,  46,  89)
NAVY_900   = RGBColor( 14,  24,  51)
GRAY_50    = RGBColor(250, 248, 245)
GRAY_900   = RGBColor( 26,  28,  38)
POS        = RGBColor(193,  59,  44)
WARN       = RGBColor(245, 158,  11)
INFO       = RGBColor( 59, 130, 246)
WHITE      = RGBColor(255, 255, 255)
TEXT_MUTE  = RGBColor(120, 110, 100)
FOOTER_CLR = RGBColor(160, 150, 140)
SUBTEXT    = RGBColor(200, 195, 190)
PROCESS_CLR= RGBColor(210, 205, 230)
ORANGE_LIGHT=RGBColor(255, 240, 225)
CODE_GREEN = RGBColor(180, 220, 180)

# ── Canvas ─────────────────────────────────────────────────────
W = Inches(20)
H = Inches(11.25)
MARGIN_X = Inches(1.396)
MARGIN_Y = Inches(0.667)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helper functions ───────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width_pt=1.0):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt2(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def add_tb(slide, x, y, w, h, text, size_pt, color_rgb,
           bold=False, italic=False, font="Noto Sans TC",
           align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color_rgb
    run.font.bold = bold
    run.font.italic = italic
    return txBox

def set_slide_bg(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_watermark(slide):
    add_tb(slide, 17.6, 10.75, 2.2, 0.35,
           "斜槓姐姐 · @byslashie", 10, TEXT_MUTE,
           font="Noto Sans TC")

def add_page_num(slide, current, total=19):
    add_tb(slide, MARGIN_X/Inches(1), 10.75, 1.5, 0.35,
           f"{current:02d} / {total:02d}", 10, TEXT_MUTE,
           font="Courier New")

def add_accent_line(slide, x, y):
    add_rect(slide, x, y, 0.8, 0.04, fill_rgb=ORANGE_500)

def add_tag(slide, x, y, text, color=None):
    c = color or ORANGE_600
    add_tb(slide, x, y, 10, 0.3, text, 13, c, bold=True)

def add_title(slide, x, y, text, size=40, color=None, bold=True):
    c = color or GRAY_900
    add_tb(slide, x, y, 17.2, 0.7, text, size, c, bold=bold)

def bullet_dot(slide, x, y, text, color=WHITE, size=18, dot_color=None):
    dc = dot_color or ORANGE_500
    add_tb(slide, x, y, 0.25, 0.35, "●", 12, dc)
    add_tb(slide, x+0.3, y, 16.5, 0.4, text, size, color)

# ═══════════════════════════════════════════════════════════════
# SLIDE 01 — 封面
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, NAVY_900)

add_rect(sl, 1.396, 2.0, 0.8, 0.04, fill_rgb=ORANGE_500)
add_tb(sl, 1.396, 2.2, 12, 0.4, "AI 策略入門班  ·  2026.05.30", 14, ORANGE_300, bold=True)
add_tb(sl, 1.396, 3.0, 14, 1.0, "不會寫程式", 72, WHITE, bold=True)
add_tb(sl, 1.396, 4.0, 14, 1.0, "也能做策略", 72, ORANGE_300, bold=True)
add_tb(sl, 1.396, 5.4, 14, 0.5, "用 AI 打造你的 XQ 投資助手", 20, SUBTEXT)
add_tb(sl, 1.396, 10.75, 4, 0.35, "byslashie.com", 11, FOOTER_CLR, font="Courier New")
add_tb(sl, 16.0,  10.75, 3.5, 0.35, "Threads @byslashie", 11, FOOTER_CLR, font="Courier New")
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 02 — 今天議程
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "TODAY'S AGENDA")
add_title(sl, 1.396, 2.0, "今天你會學到什麼")

items = [
    "以前建構策略的流程與痛點",
    "用 Perplexity + 論文找到有根據的因子",
    "Claude Skill 一步步釐清目標、建立 XS 策略",
    "XQ 回測 → 結果存進 HTML 視覺化分析",
]
for i, item in enumerate(items):
    y = 3.2 + i * 0.8
    bullet_dot(sl, 1.396, y, item, color=GRAY_900, size=20, dot_color=ORANGE_500)

add_page_num(sl, 2)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 03 — 以前怎麼做
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "BEFORE")
add_title(sl, 1.396, 2.0, "以前建構策略怎麼做？")

# First row boxes
row1 = ["有想法", "想法量化", "回測驗證"]
row2 = ["修改策略", "實戰策略", "優化策略"]

box_w = 2.0
box_h = 0.55
gap   = 2.2
y1    = 4.5
y2    = 5.8

for i, label in enumerate(row1):
    x = 1.4 + i * gap
    r = add_rect(sl, x, y1, box_w, box_h, fill_rgb=PROCESS_CLR)
    add_tb(sl, x, y1+0.05, box_w, box_h, label, 16, NAVY_700,
           bold=True, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(row1)-1:
        add_rect(sl, x+box_w, y1+0.22, gap-box_w, 0.1, fill_rgb=ORANGE_500)

# Arrow connecting row1 to row2 (vertical then horizontal)
# Simple: down arrow at end of row1
add_rect(sl, 1.4+2*gap+box_w+0.1, y1+0.22, 0.1, (y2-y1), fill_rgb=ORANGE_500)

for i, label in enumerate(row2):
    # row2 goes right to left visually, but we draw left to right in reverse
    x = 1.4 + (2-i) * gap
    r = add_rect(sl, x, y2, box_w, box_h, fill_rgb=PROCESS_CLR)
    add_tb(sl, x, y2+0.05, box_w, box_h, label, 16, NAVY_700,
           bold=True, align=PP_ALIGN.CENTER)
    if i < len(row2)-1:
        add_rect(sl, x-gap+box_w, y2+0.22, gap-box_w, 0.1, fill_rgb=ORANGE_500)

add_tb(sl, 1.396, 8.5, 17.2, 0.5,
       "耗時費力、靠感覺、沒有學術根據", 16, POS, bold=True)

add_page_num(sl, 3)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 04 — 現在的流程
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "NOW WITH AI")
add_title(sl, 1.396, 2.0, "加入 AI 之後的流程")

row1 = ["有想法", "Claude Skill", "Perplexity 論文", "Claude 寫 XS"]
row2 = ["XQ 回測", "AI 分析報告", "修改策略", "實戰"]
bw   = 1.9
bh   = 0.5
gap  = 2.2 * (17.2 / (len(row1)*2.2))   # spread across width
gap  = (17.2 - bw*len(row1)) / (len(row1)-1)

for i, label in enumerate(row1):
    x = 1.396 + i*(bw+gap)
    add_rect(sl, x, 3.5, bw, bh, fill_rgb=ORANGE_300)
    add_tb(sl, x, 3.55, bw, bh, label, 14, NAVY_900,
           bold=True, align=PP_ALIGN.CENTER)
    if i < len(row1)-1:
        add_rect(sl, x+bw, 3.72, gap, 0.08, fill_rgb=ORANGE_500)

for i, label in enumerate(row2):
    x = 1.396 + i*(bw+gap)
    add_rect(sl, x, 5.0, bw, bh, fill_rgb=ORANGE_300)
    add_tb(sl, x, 5.05, bw, bh, label, 14, NAVY_900,
           bold=True, align=PP_ALIGN.CENTER)
    if i < len(row2)-1:
        add_rect(sl, x+bw, 5.22, gap, 0.08, fill_rgb=ORANGE_500)

add_tb(sl, 1.396, 8.0, 17.2, 0.5,
       "每個環節都有 AI 幫你，省時 80%", 16, ORANGE_600, bold=True)

add_page_num(sl, 4)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 05 — 章節 01
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, NAVY_900)

add_tb(sl, 1.396, 3.0, 5, 1.2, "01", 80, ORANGE_500, bold=True, font="Courier New")
add_tb(sl, 1.396, 4.5, 14, 1.0, "為什麼從論文", 60, WHITE, bold=True)
add_tb(sl, 1.396, 5.6, 14, 1.0, "找因子？", 60, WHITE, bold=True)
add_tb(sl, 1.396, 7.0, 14, 0.4, "學術依據 vs 網路資訊", 18, TEXT_MUTE, italic=True)
add_rect(sl, 1.396, 7.6, 0.8, 0.06, fill_rgb=ORANGE_500)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 06 — 論文 vs 網路
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "THE PROBLEM")
add_title(sl, 1.396, 2.0, "為什麼要找論文？")

# Left column — gray box
lx, ly, lw, lh = 1.396, 3.0, 7.8, 5.5
add_rect(sl, lx, ly, lw, lh,
         fill_rgb=RGBColor(230,228,240),
         line_rgb=RGBColor(180,170,200), line_width_pt=1.5)
add_tb(sl, lx+0.25, ly+0.25, lw-0.3, 0.5,
       "❌  網路資訊", 22, NAVY_700, bold=True)
left_items = ["來源不明、難以驗證", "容易過時、流傳錯誤", "不知道真的有沒有效"]
for i, t in enumerate(left_items):
    bullet_dot(sl, lx+0.25, ly+1.0+i*0.9, t, color=GRAY_900, size=17, dot_color=RGBColor(180,60,60))

# Right column — orange box
rx, ry, rw, rh = 10.5, 3.0, 7.8, 5.5
add_rect(sl, rx, ry, rw, rh,
         fill_rgb=RGBColor(255,240,220),
         line_rgb=ORANGE_500, line_width_pt=1.5)
add_tb(sl, rx+0.25, ry+0.25, rw-0.3, 0.5,
       "✅  學術論文", 22, NAVY_700, bold=True)
right_items = ["實證研究、有方法論", "可重現、可驗證", "清楚的因子條件"]
for i, t in enumerate(right_items):
    bullet_dot(sl, rx+0.25, ry+1.0+i*0.9, t, color=GRAY_900, size=17, dot_color=ORANGE_600)

add_page_num(sl, 6)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 07 — 章節 02
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, NAVY_900)

add_tb(sl, 1.396, 3.0, 5, 1.2, "02", 80, ORANGE_500, bold=True, font="Courier New")
add_tb(sl, 1.396, 4.5, 14, 1.0, "找論文", 60, WHITE, bold=True)
add_tb(sl, 1.396, 5.6, 14, 1.0, "兩種方法", 60, WHITE, bold=True)
add_tb(sl, 1.396, 7.0, 14, 0.4, "免費版 & 進階版", 18, TEXT_MUTE, italic=True)
add_rect(sl, 1.396, 7.6, 0.8, 0.06, fill_rgb=ORANGE_500)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 08 — Perplexity
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "SIMPLE  ·  免費  ·  零設定")
add_title(sl, 1.396, 2.0, "Perplexity Academic 模式")

steps = [
    "開啟 perplexity.ai → 右上角切換 Academic 模式",
    "輸入英文關鍵字，例：momentum factor Taiwan stock",
    "把搜尋摘要貼給 Claude，請它分析因子條件",
]
for i, step in enumerate(steps):
    y = 3.2 + i * 1.0
    # Orange circle number
    add_rect(sl, 1.396, y, 0.4, 0.4, fill_rgb=ORANGE_500)
    add_tb(sl, 1.396, y+0.02, 0.4, 0.4,
           str(i+1), 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(sl, 1.95, y, 16.0, 0.45, step, 18, GRAY_900)

# Hint box
hx, hy, hw, hh = 1.396, 7.5, 17.2, 0.9
add_rect(sl, hx, hy, hw, hh, fill_rgb=ORANGE_LIGHT, line_rgb=ORANGE_500, line_width_pt=1.5)
add_tb(sl, hx+0.3, hy+0.15, hw-0.5, 0.6,
       "完全免費，人人可用，今天 DEMO 用這個", 18, ORANGE_600, bold=True)

add_page_num(sl, 8)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 09 — MCP 進階版
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "ADVANCED  ·  我自己用的方式")
add_title(sl, 1.396, 2.0, "MCP：Consensus + Scholar Gateway")

# Two tool cards
cards = [
    ("Consensus", "找實證結論", "確認因子是否有學術支持"),
    ("Scholar Gateway", "語意搜尋", "找最新文獻、多角度搜尋"),
]
for i, (title, sub1, sub2) in enumerate(cards):
    cx = 1.396 + i * 9.0
    cy = 3.5
    add_rect(sl, cx, cy, 8.0, 4.0, fill_rgb=WHITE, line_rgb=ORANGE_500, line_width_pt=2)
    add_tb(sl, cx+0.3, cy+0.3, 7.0, 0.6, title, 28, NAVY_700, bold=True)
    add_tb(sl, cx+0.3, cy+1.1, 7.0, 0.4, sub1, 18, ORANGE_600, bold=True)
    add_tb(sl, cx+0.3, cy+1.7, 7.0, 0.4, sub2, 16, GRAY_900)

# Bottom hint
add_rect(sl, 1.396, 8.8, 17.2, 0.7, fill_rgb=RGBColor(235,232,230))
add_tb(sl, 1.7, 8.9, 16.5, 0.5,
       "需要設定 MCP Server，適合進階使用者", 16, TEXT_MUTE, italic=True)

add_page_num(sl, 9)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — 章節 03
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, NAVY_900)

add_tb(sl, 1.396, 3.0, 5, 1.2, "03", 80, ORANGE_500, bold=True, font="Courier New")
add_tb(sl, 1.396, 4.5, 14, 1.0, "Claude Skill", 60, WHITE, bold=True)
add_tb(sl, 1.396, 5.6, 14, 1.0, "策略建立向導", 60, WHITE, bold=True)
add_tb(sl, 1.396, 7.0, 14, 0.4, "一步步釐清你的策略目標", 18, TEXT_MUTE, italic=True)
add_rect(sl, 1.396, 7.6, 0.8, 0.06, fill_rgb=ORANGE_500)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Skill 介紹
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "THE SKILL")
add_title(sl, 1.396, 2.0, "什麼是 Claude Skill？")

add_tb(sl, 1.396, 2.9, 17.2, 0.7,
       "一段 Prompt，貼進 Claude 就能用。不需要付費、不需要設定，開新對話貼上就是你的策略顧問。",
       17, GRAY_900, word_wrap=True)

use_steps = [
    "打開 claude.ai 開新對話",
    "把 Skill Prompt 貼進去",
    "開始建立你的量化策略",
]
for i, step in enumerate(use_steps):
    y = 4.1 + i * 0.75
    bullet_dot(sl, 1.396, y, step, color=GRAY_900, size=18, dot_color=ORANGE_500)

# CTA bar
add_rect(sl, 1.396, 8.0, 17.2, 0.75, fill_rgb=ORANGE_500)
add_tb(sl, 1.7, 8.1, 16.5, 0.55,
       "今天直播結束送你這個 Skill！", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_page_num(sl, 11)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Skill 五步驟
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "HOW IT WORKS")
add_title(sl, 1.396, 2.0, "Skill 一步步問你什麼")

skill_steps = [
    "策略類型（動能 / 價值 / 品質 / 低波動）",
    "策略目標（高報酬 / 高勝率 / 穩定性）",
    "方向偏好（順勢 / 逆勢）",
    "持有週期（日 / 週 / 月 / 季）",
    "目標市場（台股 / 美股 / 全球）",
]
for i, step in enumerate(skill_steps):
    y = 3.0 + i * 0.8
    add_rect(sl, 1.396, y, 0.42, 0.42, fill_rgb=ORANGE_500)
    add_tb(sl, 1.396, y+0.03, 0.42, 0.42,
           str(i+1), 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(sl, 1.95, y, 16.0, 0.45, step, 18, GRAY_900)

add_tb(sl, 1.396, 8.5, 17.2, 0.5,
       "→ 收到答案後，Skill 輸出 XS 策略框架（含基本面 + 籌碼面 + 技術面）",
       16, ORANGE_600, bold=True)

add_page_num(sl, 12)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — 章節 04
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, NAVY_900)

add_tb(sl, 1.396, 3.0, 5, 1.2, "04", 80, ORANGE_500, bold=True, font="Courier New")
add_tb(sl, 1.396, 4.5, 14, 1.0, "XQ 實作", 60, WHITE, bold=True)
add_tb(sl, 1.396, 5.6, 14, 1.0, "+ 回測", 60, WHITE, bold=True)
add_tb(sl, 1.396, 7.0, 14, 0.4, "把 AI 輸出的 XS 框架跑起來", 18, TEXT_MUTE, italic=True)
add_rect(sl, 1.396, 7.6, 0.8, 0.06, fill_rgb=ORANGE_500)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — XQ 四步驟
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "LIVE DEMO")
add_title(sl, 1.396, 2.0, "XQ 實作四步驟")

grid = [
    ("01", "選股中心設定三因子條件",  1.396, 3.2),
    ("02", "建立 XS 交易腳本（AI 生成）",  10.3,  3.2),
    ("03", "設計出場邏輯（季度換倉）", 1.396, 6.5),
    ("04", "執行回測",                10.3,  6.5),
]
for num, label, cx, cy in grid:
    add_rect(sl, cx, cy, 7.8, 3.0, fill_rgb=WHITE, line_rgb=ORANGE_500, line_width_pt=2)
    add_tb(sl, cx+0.3, cy+0.3, 2, 0.9, num, 52, ORANGE_500, bold=True, font="Courier New")
    add_tb(sl, cx+0.3, cy+1.4, 7.0, 0.5, label, 22, NAVY_700, bold=True)

add_page_num(sl, 14)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — 回測結果
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "RESULT")
add_title(sl, 1.396, 2.0, "因子 ETF 回測成績")

stat_cards = [
    ("5年最大報酬", "+112.42%", "優化版 · 3因子", POS, ORANGE_500),
    ("基礎版報酬",  "+86.42%",  "3因子基礎設定",   WARN, WARN),
    ("0050 同期",   "+62%",     "大盤對比",         TEXT_MUTE, RGBColor(160,155,150)),
]
for i, (lab, num, delta, numcol, topcol) in enumerate(stat_cards):
    cx = 1.396 + i * 5.9
    cy = 3.5
    cw, ch = 5.5, 3.5
    add_rect(sl, cx, cy, cw, ch, fill_rgb=WHITE, line_rgb=RGBColor(220,215,210), line_width_pt=1)
    # top accent
    add_rect(sl, cx, cy, cw, 0.12, fill_rgb=topcol)
    add_tb(sl, cx+0.25, cy+0.25, cw-0.3, 0.4, lab, 16, GRAY_900)
    add_tb(sl, cx+0.25, cy+0.85, cw-0.3, 1.0, num, 48, numcol, bold=True)
    add_tb(sl, cx+0.25, cy+2.7,  cw-0.3, 0.4, delta, 14, TEXT_MUTE)

add_tb(sl, 1.396, 8.2, 17.2, 0.5,
       "2022年下半年起逆轉領先大盤，2019–2022年表現弱於0050", 15, TEXT_MUTE)

add_page_num(sl, 15)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — 存進 HTML
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "ANALYSIS")
add_title(sl, 1.396, 2.0, "回測結果存進 HTML 分析")

html_steps = [
    "XQ 回測完成 → 匯出 CSV",
    "把 CSV 拖入 Demo 分析系統",
    "AI 自動計算 CAGR / Sharpe / MDD 等 KPI",
]
for i, step in enumerate(html_steps):
    y = 3.2 + i * 1.0
    bullet_dot(sl, 1.396, y, step, color=GRAY_900, size=20, dot_color=ORANGE_500)

# Hint box
add_rect(sl, 1.396, 7.0, 17.2, 0.9, fill_rgb=ORANGE_LIGHT, line_rgb=ORANGE_500, line_width_pt=1.5)
add_tb(sl, 1.7, 7.15, 16.5, 0.6,
       "視覺化報告一鍵生成，今天 DEMO 現場示範", 18, ORANGE_600, bold=True)

add_page_num(sl, 16)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — AI 分析
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, GRAY_50)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "AI REVIEW")
add_title(sl, 1.396, 2.0, "讓 AI 幫你解讀回測報告")

# Code block
add_rect(sl, 1.396, 3.0, 17.2, 4.5, fill_rgb=NAVY_900)
code_text = (
    "這是我的回測結果：[貼上數據]\n\n"
    "請幫我分析：\n"
    "1. 勝率和賠率是否合理？\n"
    "2. 這個策略的弱點在哪？\n"
    "3. 給我 3 個具體改進方向"
)
add_tb(sl, 1.7, 3.2, 16.5, 4.0, code_text, 16, CODE_GREEN, font="Courier New")

add_page_num(sl, 17)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — 禮物章節
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, NAVY_900)

add_tb(sl, 0, 2.8, 20, 1.5, "🎁", 80, WHITE, align=PP_ALIGN.CENTER)
add_tb(sl, 0, 4.3, 20, 1.0, "今天的禮物", 60, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_tb(sl, 0, 5.6, 20, 0.5, "帶走這個工具，自己找下一個因子", 20, ORANGE_300, align=PP_ALIGN.CENTER)
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — 禮物包 + 結尾
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, NAVY_900)

add_accent_line(sl, 1.396, 1.3)
add_tag(sl, 1.396, 1.5, "GIFT · 填表免費領取", color=ORANGE_300)
add_title(sl, 1.396, 2.0, "禮物包內容", size=48, color=WHITE)

gift_items = [
    "🎁  Claude Skill Prompt 文字檔（貼上即用）",
    "📝  Perplexity 論文搜尋關鍵字清單",
    "📄  XS 策略框架模板（含基本面 + 籌碼面 + 技術面）",
]
for i, item in enumerate(gift_items):
    y = 3.2 + i * 1.0
    bullet_dot(sl, 1.396, y, item, color=WHITE, size=20, dot_color=ORANGE_300)

# CTA bar
add_rect(sl, 1.396, 8.5, 17.2, 0.8, fill_rgb=ORANGE_500)
add_tb(sl, 1.7, 8.6, 16.5, 0.6,
       "直播留言「積木」→ 填表領取", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_tb(sl, 1.396, 10.75, 4, 0.35, "byslashie.com", 11, FOOTER_CLR, font="Courier New")
add_tb(sl, 16.0,  10.75, 3.5, 0.35, "Threads @byslashie", 11, FOOTER_CLR, font="Courier New")
add_watermark(sl)

# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
out_path = "/Users/chenyu/Desktop/Projects/Krystal_Core/07-工具/slides_20260530_v2.pptx"
prs.save(out_path)
print(f"✅ 投影片已生成：slides_20260530_v2.pptx （19頁）")
