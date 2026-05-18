"""
生成「不會寫程式也能做策略」AI 策略入門班投影片
斜槓姐姐品牌色系，19 頁
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ─── 品牌顏色 ──────────────────────────────────────────
ORANGE     = RGBColor(232, 115, 74)
CORAL      = RGBColor(255, 112, 150)
CYAN       = RGBColor(55,  197, 212)
NAVY       = RGBColor(14,  24,  51)
NAVY2      = RGBColor(23,  46,  89)
GRAY50     = RGBColor(250, 248, 245)
WHITE      = RGBColor(255, 255, 255)
TEXT_STRONG= RGBColor(26,  26,  26)
TEXT_MUTE  = RGBColor(120, 110, 100)
GOLD       = RGBColor(201, 145, 90)
LIGHT_ORANGE = RGBColor(255, 237, 225)
PURPLE_GRAY  = RGBColor(130, 120, 160)
RED_WARN     = RGBColor(200, 60, 60)
CODE_BG      = RGBColor(30,  30,  50)

# ─── 字體 ──────────────────────────────────────────────
FONT_ZH   = "Noto Sans TC"
FONT_NUM  = "Arial Black"

# ─── 投影片尺寸 (16:9, 33.87 × 19.05 cm) ──────────────
W = Cm(33.87)
H = Cm(19.05)

# ─── 工具函數 ──────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_name=FONT_ZH, font_size=18, bold=False,
                color=TEXT_STRONG, align=PP_ALIGN.LEFT,
                wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_name=FONT_ZH, font_size=18, bold=False,
                          color=TEXT_STRONG, align=PP_ALIGN.LEFT,
                          line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def watermark(slide):
    add_textbox(slide, "斜槓姐姐",
                W - Cm(4.5), H - Cm(1.0), Cm(4), Cm(0.8),
                font_size=9, color=TEXT_MUTE, align=PP_ALIGN.RIGHT)

def orange_label(slide, text):
    add_textbox(slide, text,
                Cm(2), Cm(1.6), Cm(20), Cm(0.8),
                font_size=13, bold=True, color=ORANGE)

def big_title(slide, text, top=Cm(3.0), color=TEXT_STRONG, font_size=36):
    add_textbox(slide, text,
                Cm(2), top, Cm(29), Cm(2.5),
                font_size=font_size, bold=True, color=color)

def section_number(slide, num_text):
    add_textbox(slide, num_text,
                Cm(2), Cm(1.5), Cm(5), Cm(2.0),
                font_name=FONT_NUM, font_size=60, bold=True, color=ORANGE)

def bottom_deco_line(slide):
    add_rect(slide, Cm(2), H - Cm(0.9), Cm(2.5), Cm(0.2), fill_color=ORANGE)

def add_bullet_dot(slide, items, left, top, gap=Cm(1.0), font_size=16):
    for i, item in enumerate(items):
        cy = top + i * gap
        # 橙點
        dot = slide.shapes.add_shape(9, left, cy + Pt(5), Cm(0.25), Cm(0.25))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE
        dot.line.fill.background()
        add_textbox(slide, item, left + Cm(0.5), cy, Cm(26), Cm(0.85),
                    font_size=font_size, color=TEXT_STRONG)

def add_circle_step(slide, items, start_left, top, circle_size=Cm(0.9), gap_x=Cm(0.4), box_w=Cm(5.5), font_size=14):
    x = start_left
    for i, (num, text) in enumerate(items):
        # 圓圈
        circ = slide.shapes.add_shape(9, x, top, circle_size, circle_size)
        circ.fill.solid()
        circ.fill.fore_color.rgb = ORANGE
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(num)
        run.font.name = FONT_NUM
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE
        # 文字
        add_textbox(slide, text, x - Cm(0.2), top + circle_size + Cm(0.15),
                    box_w, Cm(1.2), font_size=font_size, color=TEXT_STRONG,
                    align=PP_ALIGN.CENTER)
        x += box_w + gap_x

def add_flow_boxes(slide, labels, start_left, top, box_w, box_h, box_color, text_color, arrow_color, font_size=13):
    """水平流程圖 + 箭頭"""
    x = start_left
    arrow_w = Cm(0.7)
    for i, label in enumerate(labels):
        rect = add_rect(slide, x, top, box_w, box_h, fill_color=box_color)
        rect.line.fill.background()
        tf = rect.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = label
        run.font.name = FONT_ZH
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = text_color
        x += box_w
        if i < len(labels) - 1:
            # 箭頭（用填色三角近似）
            arr = slide.shapes.add_shape(
                13,  # right-arrow
                x, top + box_h * 0.3, arrow_w, box_h * 0.4
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = arrow_color
            arr.line.fill.background()
            x += arrow_w


# ══════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════

def slide01(prs):
    """封面"""
    sl = blank_slide(prs)
    set_bg(sl, NAVY)
    # 左上小標
    add_textbox(sl, "AI 策略入門班  ·  2026.05.30",
                Cm(2), Cm(1.2), Cm(20), Cm(0.8),
                font_size=13, bold=True, color=ORANGE)
    # 橙色橫線裝飾（左下）
    add_rect(sl, Cm(2), H - Cm(1.5), Cm(2.5), Cm(0.25), fill_color=ORANGE)
    # 主標
    add_textbox(sl, "不會寫程式",
                Cm(2), Cm(3.8), Cm(29), Cm(4.0),
                font_size=72, bold=True, color=WHITE)
    add_textbox(sl, "也能做策略",
                Cm(2), Cm(8.2), Cm(29), Cm(4.0),
                font_size=72, bold=True, color=WHITE)
    # 副標
    add_textbox(sl, "用 AI 打造你的 XQ 投資助手",
                Cm(2), Cm(13.2), Cm(25), Cm(1.2),
                font_size=22, color=RGBColor(200, 200, 210))
    # 右下
    add_textbox(sl, "斜槓姐姐  byslashie.com",
                W - Cm(9), H - Cm(1.2), Cm(8.5), Cm(0.8),
                font_size=11, color=RGBColor(160, 155, 170), align=PP_ALIGN.RIGHT)

def slide02(prs):
    """今天你會學到什麼"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "TODAY'S AGENDA")
    big_title(sl, "今天你會學到什麼")
    items = [
        "以前建構策略的痛點",
        "用 AI + 論文找到有根據的因子",
        "Claude Skill 一步步建立 XS 策略",
        "XQ 回測 → 結果存進 HTML 分析",
    ]
    add_bullet_dot(sl, items, Cm(2.5), Cm(6.5), gap=Cm(1.6), font_size=20)
    watermark(sl)

def slide03(prs):
    """以前怎麼做"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "BEFORE")
    big_title(sl, "以前建構策略怎麼做？")
    labels = ["有想法", "想法量化", "回測驗證", "修改策略", "實戰策略", "優化策略"]
    n = len(labels)
    total_w = Cm(29)
    box_w = Cm(3.5)
    arrow_w = Cm(0.5)
    used = box_w * n + arrow_w * (n-1)
    sx = Cm(2.5)
    add_flow_boxes(sl, labels, sx, Cm(7.5), box_w, Cm(2.2),
                   PURPLE_GRAY, WHITE, RGBColor(100, 90, 130), font_size=13)
    # 底部紅字
    add_textbox(sl, "耗時、靠感覺、不知道有沒有學術根據",
                Cm(2), Cm(14.5), Cm(29), Cm(1.0),
                font_size=15, bold=True, color=RED_WARN)
    watermark(sl)

def slide04(prs):
    """現在怎麼做"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "NOW WITH AI")
    big_title(sl, "加入 AI 之後的流程")
    labels = ["有想法", "Claude Skill\n釐清目標", "Perplexity\n找論文", "Claude\n寫 XS", "XQ\n回測驗證", "AI\n分析報告", "修改策略", "實戰"]
    n = len(labels)
    box_w = Cm(3.3)
    arrow_w = Cm(0.4)
    sx = Cm(1.2)
    add_flow_boxes(sl, labels, sx, Cm(7.0), box_w, Cm(2.6),
                   ORANGE, WHITE, RGBColor(180, 80, 30), font_size=11)
    add_textbox(sl, "每個環節都有 AI 幫你",
                Cm(2), Cm(14.5), Cm(20), Cm(1.0),
                font_size=16, bold=True, color=ORANGE)
    watermark(sl)

def slide05(prs):
    """章節切換 01"""
    sl = blank_slide(prs)
    set_bg(sl, NAVY)
    section_number(sl, "01")
    add_textbox(sl, "為什麼從論文找因子？",
                Cm(2), Cm(5.5), Cm(28), Cm(3.0),
                font_size=48, bold=True, color=WHITE)
    add_textbox(sl, "學術依據 vs 網路資訊",
                Cm(2), Cm(9.8), Cm(20), Cm(1.0),
                font_size=20, color=RGBColor(200, 200, 210))
    bottom_deco_line(sl)
    watermark(sl)

def slide06(prs):
    """論文 vs 網路"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "THE PROBLEM")
    big_title(sl, "為什麼要找論文？")
    # 左欄（灰）
    add_rect(sl, Cm(2), Cm(6.0), Cm(13), Cm(8.5), fill_color=RGBColor(220, 215, 210))
    add_textbox(sl, "網路資訊",
                Cm(2.5), Cm(6.4), Cm(12), Cm(1.2),
                font_size=20, bold=True, color=RGBColor(100, 90, 80))
    items_l = ["• 來源不明", "• 容易過時", "• 無法驗證"]
    for i, t in enumerate(items_l):
        add_textbox(sl, t, Cm(2.8), Cm(8.2) + i*Cm(1.5), Cm(11), Cm(1.0),
                    font_size=16, color=RGBColor(90,80,70))
    # 右欄（橙）
    add_rect(sl, Cm(16), Cm(6.0), Cm(15.5), Cm(8.5), fill_color=LIGHT_ORANGE)
    add_textbox(sl, "學術論文",
                Cm(16.5), Cm(6.4), Cm(14), Cm(1.2),
                font_size=20, bold=True, color=ORANGE)
    items_r = ["✓ 有實證、有樣本支撐", "✓ 有方法論可重現", "✓ 同行審查，可信度高"]
    for i, t in enumerate(items_r):
        add_textbox(sl, t, Cm(16.8), Cm(8.2) + i*Cm(1.5), Cm(14), Cm(1.0),
                    font_size=16, color=TEXT_STRONG)
    watermark(sl)

def slide07(prs):
    """章節切換 02"""
    sl = blank_slide(prs)
    set_bg(sl, NAVY)
    section_number(sl, "02")
    add_textbox(sl, "找論文兩種方法",
                Cm(2), Cm(5.5), Cm(28), Cm(3.0),
                font_size=48, bold=True, color=WHITE)
    add_textbox(sl, "免費版 & 進階版",
                Cm(2), Cm(9.8), Cm(20), Cm(1.0),
                font_size=20, color=RGBColor(200, 200, 210))
    bottom_deco_line(sl)
    watermark(sl)

def slide08(prs):
    """Perplexity 簡單版"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "SIMPLE · 免費")
    big_title(sl, "Perplexity Academic 模式")
    steps = [
        "① 開 perplexity.ai → 切換 Academic 模式",
        "② 輸入英文關鍵字（例：momentum factor Taiwan）",
        "③ 把摘要貼給 Claude 分析",
    ]
    for i, s in enumerate(steps):
        add_textbox(sl, s, Cm(2.5), Cm(6.5) + i*Cm(1.8), Cm(28), Cm(1.2),
                    font_size=18, color=TEXT_STRONG)
    # 提示框
    add_rect(sl, Cm(2), Cm(13.5), Cm(28), Cm(2.5), fill_color=LIGHT_ORANGE)
    add_textbox(sl, "完全免費，零設定，直接用",
                Cm(3), Cm(14.0), Cm(26), Cm(1.2),
                font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    watermark(sl)

def slide09(prs):
    """MCP 進階版"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "ADVANCED · 我的方式")
    big_title(sl, "MCP：Consensus + Scholar Gateway")
    add_textbox(sl, "我自己用的工具，搜尋更精準、可多角度搜尋",
                Cm(2), Cm(5.8), Cm(28), Cm(1.0),
                font_size=16, color=TEXT_MUTE)
    # 兩工具
    for i, (title, desc) in enumerate([
        ("Consensus", "找實證結論\n自動整合多篇論文觀點"),
        ("Scholar Gateway", "語意搜尋\n找最新、最相關文獻"),
    ]):
        x = Cm(2) + i * Cm(16)
        add_rect(sl, x, Cm(7.5), Cm(14.5), Cm(5.5), fill_color=WHITE)
        add_textbox(sl, title, x + Cm(0.5), Cm(8.0), Cm(13), Cm(1.2),
                    font_size=22, bold=True, color=ORANGE)
        add_textbox(sl, desc, x + Cm(0.5), Cm(9.5), Cm(13), Cm(2.5),
                    font_size=15, color=TEXT_STRONG)
    add_textbox(sl, "需要設定 MCP，進階使用者適合",
                Cm(2), Cm(14.5), Cm(28), Cm(1.0),
                font_size=14, color=TEXT_MUTE)
    watermark(sl)

def slide10(prs):
    """章節切換 03"""
    sl = blank_slide(prs)
    set_bg(sl, NAVY)
    section_number(sl, "03")
    add_textbox(sl, "Claude Skill 策略建立",
                Cm(2), Cm(5.5), Cm(28), Cm(3.0),
                font_size=48, bold=True, color=WHITE)
    add_textbox(sl, "一步步釐清你的策略目標",
                Cm(2), Cm(9.8), Cm(20), Cm(1.0),
                font_size=20, color=RGBColor(200, 200, 210))
    bottom_deco_line(sl)
    watermark(sl)

def slide11(prs):
    """Skill 介紹"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "THE SKILL")
    big_title(sl, "什麼是 Claude Skill？")
    add_textbox(sl, "一段 Prompt，貼進 Claude 就能用，完全免費，不需要 Project",
                Cm(2), Cm(6.0), Cm(29), Cm(1.2),
                font_size=18, color=TEXT_STRONG)
    add_textbox(sl, "使用方式：",
                Cm(2), Cm(8.0), Cm(10), Cm(0.8),
                font_size=16, bold=True, color=TEXT_MUTE)
    steps = ["① 開 Claude 新對話", "② 貼入 Prompt", "③ 開始建立策略"]
    for i, s in enumerate(steps):
        add_textbox(sl, s, Cm(2.5), Cm(9.0) + i*Cm(1.4), Cm(26), Cm(1.0),
                    font_size=17, color=TEXT_STRONG)
    # CTA 框
    add_rect(sl, Cm(2), Cm(14.2), Cm(29), Cm(2.4), fill_color=ORANGE)
    add_textbox(sl, "今天的禮物就是這個！",
                Cm(2), Cm(14.8), Cm(29), Cm(1.0),
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    watermark(sl)

def slide12(prs):
    """Skill 流程"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "HOW IT WORKS")
    big_title(sl, "Skill 一步步問你什麼")
    items = [
        (1, "策略類型\n動能/價值/品質/低波動"),
        (2, "策略目標\n報酬/勝率/穩定性"),
        (3, "順勢\nor 逆勢"),
        (4, "持有週期\n日/週/月/季"),
        (5, "市場範圍\n台股/美股/全球"),
    ]
    add_circle_step(sl, items, Cm(1.5), Cm(6.5), circle_size=Cm(1.0),
                    gap_x=Cm(0.3), box_w=Cm(6.0), font_size=13)
    add_textbox(sl, "→ 最後輸出 XS 策略框架",
                Cm(2), Cm(15.5), Cm(28), Cm(1.0),
                font_size=16, bold=True, color=ORANGE)
    watermark(sl)

def slide13(prs):
    """章節切換 04"""
    sl = blank_slide(prs)
    set_bg(sl, NAVY)
    section_number(sl, "04")
    add_textbox(sl, "XQ 實作 + 回測",
                Cm(2), Cm(5.5), Cm(28), Cm(3.0),
                font_size=48, bold=True, color=WHITE)
    add_textbox(sl, "把 AI 輸出的框架跑起來",
                Cm(2), Cm(9.8), Cm(20), Cm(1.0),
                font_size=20, color=RGBColor(200, 200, 210))
    bottom_deco_line(sl)
    watermark(sl)

def slide14(prs):
    """XQ 四步驟"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "LIVE DEMO")
    big_title(sl, "XQ 四步驟實作")
    steps = [
        ("01", "選股中心設定三因子條件"),
        ("02", "建立交易腳本（XS 貼入）"),
        ("03", "設計出場邏輯（季度換倉）"),
        ("04", "執行回測"),
    ]
    for i, (num, desc) in enumerate(steps):
        x = Cm(1.5) + i * Cm(8.0)
        add_textbox(sl, num, x, Cm(6.5), Cm(7.5), Cm(2.8),
                    font_name=FONT_NUM, font_size=42, bold=True, color=ORANGE)
        add_textbox(sl, desc, x, Cm(9.5), Cm(7.8), Cm(2.5),
                    font_size=15, color=TEXT_STRONG)
    watermark(sl)

def slide15(prs):
    """回測結果"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "RESULT")
    big_title(sl, "回測結果")
    cards = [
        ("+112.42%", "5年最大報酬", ORANGE, WHITE),
        ("+86.42%",  "基礎版報酬",  NAVY2,  WHITE),
        ("0050 +62%","同期大盤",    RGBColor(200,195,190), TEXT_STRONG),
    ]
    for i, (num, label, bg, fg) in enumerate(cards):
        x = Cm(2) + i * Cm(10.5)
        add_rect(sl, x, Cm(6.0), Cm(10.0), Cm(6.5), fill_color=bg)
        add_textbox(sl, num, x, Cm(6.8), Cm(10), Cm(3.5),
                    font_name=FONT_NUM, font_size=38, bold=True, color=fg,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, label, x, Cm(10.0), Cm(10), Cm(1.2),
                    font_size=15, color=fg, align=PP_ALIGN.CENTER)
    add_textbox(sl, "2022年下半年起逆轉領先大盤",
                Cm(2), Cm(15.5), Cm(28), Cm(1.0),
                font_size=14, color=TEXT_MUTE)
    watermark(sl)

def slide16(prs):
    """存入 HTML"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "ANALYSIS")
    big_title(sl, "回測結果存進 HTML 分析")
    steps = [
        "① XQ 匯出回測 CSV",
        "② 拖入 demo 分析系統",
        "③ AI 自動計算 CAGR / Sharpe / MDD",
    ]
    for i, s in enumerate(steps):
        add_textbox(sl, s, Cm(2.5), Cm(6.8) + i*Cm(2.0), Cm(28), Cm(1.2),
                    font_size=19, color=TEXT_STRONG)
    add_rect(sl, Cm(2), Cm(14.0), Cm(29), Cm(2.5), fill_color=LIGHT_ORANGE)
    add_textbox(sl, "視覺化報告，一鍵生成",
                Cm(2), Cm(14.6), Cm(29), Cm(1.0),
                font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    watermark(sl)

def slide17(prs):
    """AI 分析"""
    sl = blank_slide(prs)
    set_bg(sl, GRAY50)
    orange_label(sl, "AI REVIEW")
    big_title(sl, "讓 AI 幫你解讀報告")
    # code block 背景
    add_rect(sl, Cm(2), Cm(6.2), Cm(29), Cm(9.5), fill_color=CODE_BG)
    code_lines = [
        "這是我的回測結果：[貼數據]",
        "請幫我分析：",
        "",
        "1. 勝率和賠率是否合理？",
        "2. 策略的弱點在哪？",
        "3. 給我 3 個改進方向",
    ]
    add_multiline_textbox(sl, code_lines, Cm(2.8), Cm(6.8), Cm(27), Cm(8.5),
                          font_name="Courier New", font_size=17,
                          color=RGBColor(180, 230, 180))
    watermark(sl)

def slide18(prs):
    """章節切換：禮物"""
    sl = blank_slide(prs)
    set_bg(sl, NAVY)
    add_textbox(sl, "今天的禮物",
                Cm(2), Cm(4.5), Cm(29), Cm(3.5),
                font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "帶走這個工具，自己找下一個因子",
                Cm(2), Cm(10.5), Cm(29), Cm(1.5),
                font_size=22, color=RGBColor(200, 200, 210), align=PP_ALIGN.CENTER)
    bottom_deco_line(sl)
    watermark(sl)

def slide19(prs):
    """禮物包 + 結尾"""
    sl = blank_slide(prs)
    set_bg(sl, NAVY)
    add_textbox(sl, "GIFT",
                Cm(2), Cm(1.3), Cm(20), Cm(0.8),
                font_size=13, bold=True, color=ORANGE)
    add_textbox(sl, "禮物包內容",
                Cm(2), Cm(2.5), Cm(29), Cm(2.0),
                font_size=40, bold=True, color=WHITE)
    gifts = [
        "🎁  Claude Skill Prompt 文字檔（貼上即用）",
        "📝  Perplexity 論文搜尋關鍵字清單",
        "📄  XS 策略框架模板",
    ]
    for i, g in enumerate(gifts):
        add_textbox(sl, g, Cm(2.5), Cm(6.0) + i*Cm(2.0), Cm(28), Cm(1.4),
                    font_size=20, color=WHITE)
    # CTA
    add_rect(sl, Cm(2), Cm(13.5), Cm(29), Cm(2.0), fill_color=ORANGE)
    add_textbox(sl, "填表領取 → byslashie.com",
                Cm(2), Cm(13.9), Cm(29), Cm(1.0),
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "斜槓姐姐  ·  下次見！",
                Cm(2), Cm(16.8), Cm(29), Cm(1.0),
                font_size=16, color=RGBColor(180, 175, 190), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    slide01(prs)
    slide02(prs)
    slide03(prs)
    slide04(prs)
    slide05(prs)
    slide06(prs)
    slide07(prs)
    slide08(prs)
    slide09(prs)
    slide10(prs)
    slide11(prs)
    slide12(prs)
    slide13(prs)
    slide14(prs)
    slide15(prs)
    slide16(prs)
    slide17(prs)
    slide18(prs)
    slide19(prs)

    out = "/Users/chenyu/Desktop/Projects/Krystal_Core/07-工具/slides_20260530.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()
